### This code is property of the GGAO ###


# Native imports
import os
import json
import string
import logging
from typing import Tuple
from io import BytesIO
from shutil import copyfile, copy2
from collections import Counter
from abc import abstractmethod

# Installed imports
import xlrd
import magic
import zipfile
import pandas as pd
import xml.etree.ElementTree as et
# from google.oauth2 import service_account
from docx import Document
from docx.document import Document as doctwo
from docx.oxml.table import CT_Tbl
from docx.oxml.text.paragraph import CT_P
from docx.text.paragraph import Paragraph
from docx.table import _Cell, Table
from pptx import Presentation
from openpyxl import load_workbook
from pdfminer.converter import TextConverter
from pdfminer.layout import LAParams, LTTextBoxHorizontal, LTChar, LTAnno
from pdfminer.pdfdocument import PDFDocument
from pdfminer.pdfinterp import PDFResourceManager, PDFPageInterpreter
from pdfminer.pdfinterp import resolve1
from pdfminer.pdfpage import PDFPage
from pdfminer.pdfparser import PDFParser
from pdfminer.converter import PDFPageAggregator
from wand.color import Color
from wand.image import Image

# Custom imports
from genai_sdk_services.resources.vision.BaseOCR import  Textract, TesseractOCR, BaseOCR, FormRecognizer, LLMOCR
from genai_sdk_services.resources.vision.ocr2visionfeatures import detect_document, runOCR, get_blocks_cells


FULL_TEXT_KEY = "text"


class BaseFileService:
    TYPE_FILE = None

    @abstractmethod
    def get_text(self, file):
        """ Get text from file

        :param file: File to get text from
        :return: (string) Text of the file
        """
        pass

    @abstractmethod
    def get_multiple_text(self, files):
        """ Get text from file

        :param files: Files to get text from
        :return: (string) Text of the file
        """
        pass

    def get_text_from_bytes(self, file):
        """ Get text from file

        :param file: File to get text from
        :return: (string) Text of the file
        """
        raise NotImplementedError("This method is not implemented yet")

    @abstractmethod
    def extract_images(self, file, **kwargs):
        """ Extract texts from file and store it into a text file

        :param file: File to extract text from
        """
        pass

    @abstractmethod
    def extract_text(self, file):
        """ Extract text from file and store it into a text file

        :param file: File to extract text from
        """
        pass

    @abstractmethod
    def extract_multiple_text(self, files):
        """ Extract text from several files and store it into their corresponding text file

        :param files: Files to extract text from
        """
        pass

    @abstractmethod
    def get_number_pages(self, file):
        """ Extract text from several files and store it into their corresponding text file

        :param file: Files to extract text from
        """
        pass

    @classmethod
    def check_file(cls, type_file):
        """ Check if it is a valid type for the file

        :param type_file: Type to check
        :return: (bool) True if the type is valid
        """
        return type_file in cls.TYPE_FILE


class ImageService(BaseFileService):
    TYPE_FILE = ["jpeg", "jpg", "png", "svg", "tiff", "ps"]
    OCR_SERVICES = [Textract(), TesseractOCR(), FormRecognizer(), LLMOCR()]

    ocr_credentials = {}

    def __init__(self):
        self.logger = logging.getLogger(__name__)

    def get_text_from_bytes(self, file: str):
        """ Get text from file

        :param file: File to get text from
        :return: (string) Text of the file
        """
        raise NotImplementedError("This method is not implemented yet")

    def get_text(self, file: str, **kwargs: dict) -> Tuple[str, str, str, str, str]:
        """ Get text from file

        :param file: File to get text from
        :return: (tuple) Texts of the file
        """
        self.logger.debug("Running OCR")

        if "ocr_origin" in kwargs:
            # Setting OCR credentials
            if "ocr_credentials" in kwargs:
                ocr_origin = kwargs["ocr_origin"]
                ocr_credentials = kwargs["ocr_credentials"]
                if ocr_origin not in self.ocr_credentials:
                    self.ocr_credentials[ocr_origin] = ocr_credentials
                ocr_credential = service_account.Credentials.from_service_account_info(self.ocr_credentials[ocr_origin])

                # Running the OCR
                ocr_infos = runOCR(ocr_credential, [file])
                extract_blocks = kwargs.get('extract_blocks', True)
                documents_blocks = {}
                documents_paragraphs = {}
                documents_words = {}
                documents_lines = {}
                if extract_blocks:
                    documents_blocks, documents_paragraphs, documents_words, documents_lines = get_blocks_cells(ocr_infos)

                returning_texts = []
                for ocr_info in ocr_infos:
                    text = ocr_info[0].text
                    returning_texts.append(text)

                return returning_texts[0], documents_blocks[0], documents_paragraphs[0], documents_words[0], documents_lines[0]

            else:
                raise ValueError("'ocr_credentials' is not defined in kwargs")

        else:
            raise ValueError("'ocr_origin' is not defined in kwargs")

    def get_multiple_text(self, files: list, **kwargs: dict) -> Tuple[list, list, list, list, list, list]:
        """ Get text from file

        :param files: (list) Files to get text from
        :return: (tuple) Texts of the file
        """
        self.logger.debug("Running OCR")

        if "ocr_origin" in kwargs:
            # Setting OCR credentials
            ocr_origin = kwargs['ocr_origin']
            credentials = kwargs.get('ocr_credentials', {})

            if "bucket" in kwargs:
                credentials.update({'bucket': kwargs['bucket']})

            ocr_service = None
            for service in self.OCR_SERVICES:
                if service.check_origin(ocr_origin):
                    ocr_service = service
            assert isinstance(ocr_service, BaseOCR)

            returning_texts, documents_blocks, documents_paragraphs, documents_words, returning_tables, documents_lines = ocr_service.run_ocr(credentials, files, **kwargs)
        else:
            raise ValueError("'ocr_origin' is not defined in kwargs")

        return returning_texts, documents_blocks, documents_paragraphs, documents_words, returning_tables, documents_lines

    def get_multiple_text_for_extraction(self, files: list, **kwargs: dict) -> list:
        """ Get text from file

        :param files: (list) File to get text from
        :param kwargs: (dict) Arguments for the OCR
        :return: (string) Text of the file
        """
        self.logger.debug("Running OCR")

        if "ocr_origin" in kwargs:
            # Setting OCR credentials
            if "ocr_credentials" in kwargs:
                ocr_origin = kwargs["ocr_origin"]
                ocr_credentials = kwargs["ocr_credentials"]
                if ocr_origin not in self.ocr_credentials:
                    self.ocr_credentials[ocr_origin] = ocr_credentials
                ocr_credential = service_account.Credentials.from_service_account_info(self.ocr_credentials[ocr_origin])

                # Running the OCR
                ocr_infos = detect_document(
                    ocr_credential,
                    files,
                    expansion=1 / 6,
                    rotated=False,
                    distorsion_angle=10000,
                )

                returning_texts = []

                for ocr_info in ocr_infos:
                    blocks = ocr_info[0]

                    text = ""
                    for block in blocks:
                        line = block[4]
                        text += line + "\n"

                    returning_texts.append(text)

                return returning_texts
            else:
                raise ValueError("'ocr_credentials' is not defined in kwargs")
        else:
            raise ValueError("'ocr_origin' is not defined in kwargs")

    def extract_images(self, file: str, **kwargs: dict) -> list:
        """ Extract images from file and store them

        :param file: (str) File to extract images from
        :param kwargs: (dict) Arguments for the extraction
        :return: (list) List of images extracted
        """
        if not os.path.exists("images/"):
            os.mkdir("images")
        filename = "images/" + os.path.splitext(os.path.basename(file))[0] + "_pag_0.jpeg"

        copyfile(file, filename)

        return [{
            'filename': filename,
            'number': 0
        }]

    def extract_text(self, file: str, **kwargs: dict) -> Tuple[list, list, list, list, list]:
        """ Extract text from file and store it into a text file

        :param file: (str) File to extract text from
        :param kwargs: (dict) Arguments for the extraction
        :return: (tuple) Texts of the file
        """
        text, document_blocks, document_paragraphs, document_words, document_lines = self.get_text(file, **kwargs)

        text_files = []
        block_files = []
        paragraphs_files = []
        words_files = []
        lines_files = []

        fn = os.path.splitext(file)[0]
        with open(fn + ".txt", "+w", encoding='utf8') as txt:
            txt.write(text)
        with open(fn + ".json", "+w", encoding='utf8') as f:
            json.dump(document_blocks, f)
        with open(fn + "_paragraphs.json", "+w", encoding='utf8') as f:
            json.dump(document_paragraphs, f)
        with open(fn + "_words.json", "+w", encoding='utf8') as f:
            json.dump(document_words, f)
        with open(fn + "_lines.json", "+w", encoding='utf8') as f:
            json.dump(document_lines, f)

        text_files.append(fn + ".txt")
        block_files.append(fn + ".json")
        paragraphs_files.append(fn + "_paragraphs.json")
        words_files.append(fn + "_words.json")
        lines_files.append(fn + "_lines.json")

        return text_files, block_files, paragraphs_files, words_files, lines_files

    def extract_multiple_text(self, files: list, **kwargs: dict) -> Tuple[list, list, list, list, list, list]:
        """ Extract texts from file and store it into a texts file

        :param files: (list) Files to extract texts from
        :param kwargs: (dict) Arguments for the extraction
        :return: (tuple) Texts of the file
        """
        texts, documents_blocks, documents_paragraphs, documents_words, tables, lines = self.get_multiple_text(files, **kwargs)

        text_files = []
        block_files = []
        paragraphs_files = []
        words_files = []
        tables_files = []
        lines_files = []
        l_tmp = [files, texts, documents_blocks, documents_paragraphs, documents_words, tables, lines]
        n_files = len(files)
        no_empty = [i if len(i) > 0 else [[]] * n_files for i in l_tmp]

        for file in files:
            if not os.path.exists(os.path.dirname(file)):
                os.makedirs(os.path.dirname(file), exist_ok=True)

        for file, text, document_blocks, document_paragraphs, document_words, table, lines in \
                zip(*no_empty):
            fn = os.path.splitext(file)[0]
            with open(fn + ".txt", "+w", encoding='utf8') as txt:
                txt.write(text)
            with open(fn + ".json", "+w", encoding='utf8') as f:
                json.dump(document_blocks, f)
            with open(fn + "_paragraphs.json", "+w", encoding='utf8') as f:
                json.dump(document_paragraphs, f)
            with open(fn + "_words.json", "+w", encoding='utf8') as f:
                json.dump(document_words, f)

            text_files.append(fn + ".txt")
            block_files.append(fn + ".json")
            paragraphs_files.append(fn + "_paragraphs.json")
            words_files.append(fn + "_words.json")

            if kwargs.get('extract_tables', False):
                with open(fn + "_tables.json", "+w", encoding='utf8') as f:
                    json.dump(table, f)
                tables_files.append(fn + "_tables.json")
            if kwargs.get('do_lines', False):
                with open(fn + "_lines.json", "+w", encoding='utf8') as f:
                    json.dump(lines, f)
                lines_files.append(fn + "_lines.json")

        return text_files, block_files, paragraphs_files, words_files, tables_files, lines_files

    def get_number_pages(self, file: str) -> int:
        """ Get the number of pages of the file

        :param file: (str) File to extract text from
        :return: (int) Number of pages of the file
        """
        raise NotImplementedError("This method is not implemented yet")


class PdfService(BaseFileService):
    TYPE_FILE = ["pdf"]

    # constructor
    def __init__(self):
        self.logger = logging.getLogger(__name__)

    @staticmethod
    def _get_most_common(input_list: list) -> list:
        """ Get the most common element in a list

        :param input_list: (list) List to get the most common element from
        :return: (list) List of the most common element
        """
        return sorted(Counter(input_list).items(), key=lambda x: x[1], reverse=True)[0][0] if input_list else None

    def parse_element(self, element: dict, bboxes: list, cells: list, lines: list, text_tmp: str, coords_tmp: tuple, do_lines: bool, **metadata: dict) -> Tuple[list, list, list, str, tuple, dict]:
        """ Parse the element of the PDF

        :param element: (dict) Element to parse
        :param bboxes: (list) Bounding boxes of the element
        :param cells: (list) Cells of the element
        :param lines: (list) Lines of the element
        :param text_tmp: (str) Temporary text
        :param coords_tmp: (tuple) Temporary coordinates
        :param do_lines: (bool) Do lines
        :param metadata: (dict) Metadata of the element
        :return: (tuple) Bounding boxes, cells, lines, text, coordinates, metadata
        """
        if isinstance(element, LTTextBoxHorizontal):
            text_block = element.get_text()
            c0_block, r0_block, c1_block, r1_block = element.bbox
            bbox = {
                'r0': r0_block,
                'c0': c0_block,
                'r1': r1_block,
                'c1': c1_block,
                'text': text_block,
                'confidence': 1
            }
            sizes, fonts = [], []
            for line in element:
                sizes_line, fonts_line = [], []
                text_line = line.get_text()
                c0_line, r0_line, c1_line, r1_line = line.bbox
                for char in line:
                    if isinstance(char, LTChar):
                        sizes_line.append(char.size)
                        fonts_line.append(char.fontname)
                sizes.extend(sizes_line)
                fonts.extend(fonts_line)

                if do_lines:
                    metadata_line = {
                        'r0': r0_line,
                        'c0': c0_line,
                        'r1': r1_line,
                        'c1': c1_line,
                        'text': text_line,
                        'confidence': 1
                    }
                    if sizes:
                        font_line = self._get_most_common(fonts_line)
                        size_line = self._get_most_common(sizes_line)

                        metadata_line['bold'] = "bold" in font_line.lower()
                        metadata_line['italic'] = "italic" in font_line.lower()
                        metadata_line['font'] = font_line
                        metadata_line['fontsize'] = size_line
                    lines.append(metadata_line)
            metadata = {}
            if sizes:
                font = self._get_most_common(fonts)
                size = self._get_most_common(sizes)

                metadata['bold'] = "bold" in font.lower()
                metadata['italic'] = "italic" in font.lower()
                metadata['font'] = font
                metadata['fontsize'] = size

            bbox.update(metadata)
            bboxes.append(bbox)

        if isinstance(element, LTAnno) or (isinstance(element, LTChar) and element.get_text() in string.whitespace):
            r0, c0, r1, c1 = coords_tmp
            cells.append({
                'r0': r0,
                'c0': c0,
                'r1': r1,
                'c1': c1,
                'text': text_tmp,
                'font': metadata.get('font', ""),
                'bold': metadata.get('bold', False),
                'italic': metadata.get('italic', False),
                'fontsize': metadata.get('fontsize', 0),
                'confidence': 1
            })
            return bboxes, cells, lines, "", [10 ** 9, 10 ** 9, 0, 0], {}
        elif isinstance(element, LTChar):
            metadata['bold'] = metadata.get('bold', False) or "bold" in element.fontname.lower()
            metadata['italic'] = metadata.get('italic', False) or "italic" in element.fontname.lower()
            metadata['font'] = metadata['font'] if metadata.get('font', "") else element.fontname
            metadata['fontsize'] = metadata['fontsize'] if metadata.get('fontsize', "") else element.size

            text_tmp += element.get_text()

            c0_tmp, r0_tmp, c1_tmp, r1_tmp = element.bbox
            coords_tmp[0] = min(r0_tmp, coords_tmp[0])
            coords_tmp[1] = min(c0_tmp, coords_tmp[1])
            coords_tmp[2] = max(r1_tmp, coords_tmp[2])
            coords_tmp[3] = max(c1_tmp, coords_tmp[3])
            return bboxes, cells, lines, text_tmp, coords_tmp, metadata
        else:
            try:
                for el in element:
                    bboxes, cells, lines, text_tmp, coords_tmp, metadata = self.parse_element(el, bboxes, cells, lines,
                                                                                              text_tmp, coords_tmp, do_lines, **metadata)
            except Exception:
                pass

        return bboxes, cells, lines, text_tmp, coords_tmp, metadata

    def get_text(self, file: str, **kwargs: dict) -> Tuple[dict, list, list, list]:
        """ Get text from pdf fle

        :param file: (str) PDF file to get text from
        :param kwargs: (dict) Additional parameters
        :return: (tuple) Text of the pdf file
        """
        with open(file, "rb") as f:
            return self.get_text_from_bytes(f, **kwargs)

    def get_multiple_text(self, files: list, **kwargs: dict):
        """ Get text from multiple txt files

        :param files: (list) List of txt files to get text from
        :param kwargs: (dict) Additional parameters
        """
        raise NotImplementedError("This method is not implemented yet")

    def get_text_from_bytes(self, file: str, **kwargs: dict) -> Tuple[dict, list, list, list]:
        """ Get text from pdf fle

        :param file: (str) PDF file to get text from
        :param kwargs: (dict) Additional parameters
        :return: (tuple) dictionary with text separated into pages and severals list with text, cells and lines
        """
        laparams = LAParams()
        laparams_tmp = kwargs.get("laparams", None)
        if type(laparams_tmp) is float or type(laparams_tmp) is int:
            laparams = LAParams(boxes_flow=laparams_tmp)

        # PDFMiner boilerplate
        rsrcmgr = PDFResourceManager()
        bio = BytesIO()
        device = TextConverter(rsrcmgr, bio, laparams=laparams)
        interpreter = PDFPageInterpreter(rsrcmgr, device)

        # Extract text
        parser = PDFParser(file)
        document = PDFDocument(parser)

        num_pag_ini = kwargs.get("num_pag_ini", 0)
        page_limit = kwargs.get("page_limit", 10)

        extraction = {'text': ""}
        for i, page in enumerate(PDFPage.create_pages(document)):
            if num_pag_ini <= i < page_limit:
                interpreter.process_page(page)
                text = bio.getvalue()
                text = text.decode("utf-8", "replace")
                extraction[f"pag_{i}"] = text
                extraction[FULL_TEXT_KEY] += text
                bio.truncate()
                bio.seek(0)
        bio.close()
        device.close()

        pages_cells = []
        pages_blocks = []
        pages_lines = []
        if kwargs.get('do_cells', True):
            device = PDFPageAggregator(rsrcmgr, laparams=laparams)
            interpreter_blocks = PDFPageInterpreter(rsrcmgr, device)
            pages = list(PDFPage.get_pages(file, check_extractable=False))

            for n_page, page in enumerate(pages):
                _, _, x_max, y_max = page.mediabox
                interpreter_blocks.process_page(page)
                # receive the LTPage object for the page.
                layout = device.get_result()
                layout_dict = {
                    'x_max': x_max,
                    'y_max': y_max,
                    'layout': layout
                }

                cells = []
                bboxes = []
                lines = []

                # If rotation x and y are inverted in our coordinates system (assumed only multiples of 90º)
                x_max = layout_dict['x_max'] if page.rotate not in [90, 270] else layout_dict['y_max']
                y_max = layout_dict['y_max'] if page.rotate not in [90, 270] else layout_dict['x_max']
                layout = layout_dict['layout']
                for element in layout:
                    try:
                        text = ""
                        r0, c0, r1, c1 = 10 ** 9, 10 ** 9, 0, 0
                        bboxes, cells, lines, text_tmp, coords_tmp, metadata = self.parse_element(element, bboxes, cells, lines, text,
                                                                                                  [r0, c0, r1, c1], kwargs.get('do_lines', False))
                        for bbox in bboxes:
                            bbox['x_max'] = x_max
                            bbox['y_max'] = y_max
                            bbox['rotation'] = page.rotate
                            bbox['page'] = n_page
                        for cell in cells:
                            cell['x_max'] = x_max
                            cell['y_max'] = y_max
                            cell['rotation'] = page.rotate
                            cell['page'] = n_page
                        for line in lines:
                            line['x_max'] = x_max
                            line['y_max'] = y_max
                            line['rotation'] = page.rotate
                            line['page'] = n_page
                    except Exception as e:
                        self.logger.error(f"Pdf parsing failed: {e}", exc_info=True)
                pages_cells.append(cells)
                pages_blocks.append(bboxes)
                pages_lines.append(lines)

        return extraction, pages_blocks, pages_cells, pages_lines

    def extract_images(self, filename: str, **kwargs: dict) -> list:
        """ Extract the images of the file and save them

        :param filename: File to extract the images from
        :param kwargs: Additional parameters
        :return: List of images extracted
        """
        with open(filename, "rb") as f:
            file = f.read()
            f.seek(0)
            total_pages = self._extract_number_pages(f)

        num_pag_ini = kwargs.get('num_pag_ini', 0)
        page_limit = kwargs.get('page_limit', 10)
        images = []

        for page in range(num_pag_ini, min(page_limit, total_pages)):
            if int(page) < total_pages:
                with Image(blob=file, resolution=200, format="PDF[" + str(page) + "]") as document:
                    filename_tmp = os.path.splitext(filename)[0] + f"_pag_{page}.jpeg"
                    image = Image(document.sequence[0])

                    # We set the background to the image to white, to read it correctly.
                    image.background_color = Color("white")
                    image.alpha_channel = "remove"

                    # All the images will be save in jpg format.
                    image.format = "jpeg"
                    # We save the images in a temporal folder.
                    image.save(filename=filename_tmp)
                    images.append({"filename": filename_tmp, "number": int(page)})
                    # images.append({"filename": filename_tmp, "image": image.make_blob(), "number": int(page)})

        return images

    def extract_text(self, file: str, **kwargs: dict) -> Tuple[str, list, list, list, list]:
        """ Extract text from file and store it into a text file

        :param file: (str) File to extract text from
        :param kwargs: Additional parameters
        :return: (tuple) Filename, list of name of files, blocks, words and lines
        """
        extraction, blocks, words, lines = self.get_text(file, **kwargs)
        filename = ""
        filenames_pages = []
        for key in extraction:
            fn = os.path.splitext(file)[0]
            if key == FULL_TEXT_KEY:
                fn += ".txt"
                filename = fn
            else:
                fn += f"_{key}.txt"
                filenames_pages.append(fn)

            with open(fn, "+w") as f:
                for c in extraction[key]:
                    try:
                        f.write(c)
                    except:
                        continue

        filename_blocks = os.path.splitext(file)[0] + f"_blocks.json"
        with open(filename_blocks, "+w", encoding='utf8') as f:
            json.dump(blocks, f)
        filename_words = os.path.splitext(file)[0] + f"_words.json"
        with open(filename_words, "+w", encoding='utf8') as f:
            json.dump(words, f)
        filename_lines = os.path.splitext(file)[0] + f"_lines.json"
        with open(filename_lines, "+w", encoding='utf8') as f:
            json.dump(lines, f)

        return filename, filenames_pages, filename_blocks, filename_words, filename_lines

    def extract_multiple_text(self, files: list) -> Tuple[list, list, list, list, list]:
        """ Extract text from multiple files

        :param files: (list) List of files to extract text from
        :return: (tuple) List of filenames, list of names of files, blocks, words and lines
        """
        filenames = []
        filenames_pages = []
        filenames_blocks = []
        filenames_words = []
        filenames_lines = []
        for file in files:
            filename_tmp, filenames_pages_tmp, filename_blocks_tmp, filename_words_tmp, filename_lines_tmp = self.extract_text(file)
            filenames.append(filename_tmp)
            filenames_pages.append(filenames_pages_tmp)
            filenames_blocks.append(filename_blocks_tmp)
            filenames_words.append(filename_words_tmp)
            filenames_lines.append(filename_lines_tmp)

        return filenames, filenames_pages, filenames_blocks, filenames_words, filenames_lines

    def get_number_pages(self, filename: str) -> int:
        """ Get the number of pages of the pdf file

        :param filename: (str) File to get the number of pages from
        :return: (int) Number of pages of the pdf file
        """
        with open(filename, "rb") as f:
            f.seek(0)
            return self._extract_number_pages(f)

    @staticmethod
    def _extract_number_pages(file: str) -> int:
        """ Extract the number of pages of the pdf file

        :param file: (str) File to extract the number of pages from
        :return: (int) Number of pages of the pdf file
        """
        parser = PDFParser(file)
        document = PDFDocument(parser)
        pages = resolve1(document.catalog["Pages"])["Count"]
        while type(pages) != int:
            pages = resolve1(pages)
            if pages is None:
                pages = len(list(PDFPage.get_pages(file)))

        return pages


class TxtService(BaseFileService):
    TYPE_FILE = ["txt", "plain"]

    def __init__(self):
        self.logger = logging.getLogger(__name__)

    def get_text(self, file: str, **kwargs: dict) -> Tuple[dict, list, list, list]:
        """ Get text from txt fle

        :param file: (str) txt file to get text from
        :return: (tuple) Text of the txt file and empty lists
        """
        with open(file, "r", encoding="utf-8") as f:
            text = f.read()
        extraction = {FULL_TEXT_KEY: text}

        return extraction, [], [], []

    def get_multiple_text(self, files: list, **kwargs: dict):
        """ Get text from multiple txt files

        :param files: (list) List of txt files to get text from
        :param kwargs: (dict) Additional parameters
        """
        raise NotImplementedError("This method is not implemented yet")

    def get_text_from_bytes(self, file: str, **kwargs: dict) -> str:
        """ Get text from txt fle

        :param file: (str) txt file to get text from
        :return: (str) Text of the txt file
        """
        return file.decode()

    def extract_images(self, file: str, **kwargs: dict) -> list:
        """ Extract the images of the file and save them

        :param file: (str) File to extract the images from
        :param kwargs: (dict) Additional parameters
        :return: (list) List of images extracted
        """
        raise NotImplementedError("This method is not implemented yet")

    def extract_text(self, file: str, **kwargs: dict) -> Tuple[str, list, list, list]:
        """ Extract text from file and store it into a text file

        :param file: (str) File to extract text from
        :param kwargs: (dict) Additional parameters
        :return: (tuple) Filename, list of name of files, blocks, words and lines
        """
        extraction = self.get_text(file)
        text = extraction[0]['text']
        cells = extraction[1]
        words = extraction[2]

        filename = os.path.splitext(file)[0] + ".txt"
        with open(filename, "+w") as f:
            f.write(text)

        filename_pages = os.path.splitext(file)[0] + "_pag_0.txt"
        with open(filename_pages, "+w") as f:
            f.write(text)

        filename_blocks = os.path.splitext(file)[0] + f"_blocks.json"
        with open(filename_blocks, "+w", encoding='utf8') as f:
            json.dump(cells, f)
        filename_words = os.path.splitext(file)[0] + f"_words.json"
        with open(filename_words, "+w", encoding='utf8') as f:
            json.dump(words, f)

        return filename, filename_pages, filename_blocks, filename_words

    @staticmethod
    def write(text: str, filename: str):
        """ Write text into a file

        :param text: (str) Text to write into the file
        :param filename: (str) File to write the text into
        """
        with open(filename, "w") as f:
            f.write(text)

    def extract_multiple_text(self, files: list) -> list:
        """ Extract text from multiple files

        :param files: (list) List of files to extract text from
        :return: (list) List of filenames"""
        filenames = []
        for file in files:
            filenames.append(self.extract_text(file))

        return filenames

    def get_number_pages(self, file: str) -> int:
        """ Get the number of pages of the txt file

        :param file: (str) File to get the number of pages from
        :return: (int) Number of pages of the txt file
        """
        raise NotImplementedError("TXTs don't have pages.")


class WordService(BaseFileService):
    TYPE_FILE = ["docx"]

    def __init__(self):
        self.logger = logging.getLogger(__name__)

    @staticmethod
    def iter_block_items(parent: dict):
        """Yield each paragraph and table child within *parent*, in document order.
        Each returned value is an instance of either Table or Paragraph. *parent*
        would most commonly be a reference to a main Document object, but
        also works for a _Cell object, which itself can contain paragraphs and tables.

        :param parent: object to process
        """
        if isinstance(parent, doctwo):
            parent_elm = parent.element.body
        elif isinstance(parent, _Cell):
            parent_elm = parent._tc
        else:
            raise ValueError("something's not right")

        for child in parent_elm.iterchildren():
            if isinstance(child, CT_P):
                yield Paragraph(child, parent)
            elif isinstance(child, CT_Tbl):
                yield Table(child, parent)

    @staticmethod
    def get_style_paragraph(paragraph: dict) -> dict:
        """ Get the style of the paragraph

        :param paragraph: (dict) Paragraph to get the style from
        :return: (dict) Style of the paragraph
        """
        font = paragraph.style.font
        style = {
            'text': paragraph.text,
            'bold': True if font.bold is not None else False,
            'italic': True if font.italic is not None else False,
            'underline': True if font.underline is not None else False,
            'style': paragraph.style.name,
            'font': font.name if font.name else "",
            'fontSize': font.size.pt if font.size else "",
            'page': 0
        }
        return style

    def get_text(self, file: str, **kwargs: dict) -> Tuple[dict, list, list, list]:
        """ Get text from txt fle

        :param file: txt file to get text from
        :param kwargs: Additional parameters
        :return: (tuple) Dictionary with text complete and by page and list of blocks, words and lines
        """
        document = Document(file)
        text = ""
        document_paragraphs = []
        document_tables = []
        document_lines = []

        for block in self.iter_block_items(document):
            if type(block) is Paragraph:
                text += block.text + "\n"
                if len(block.text) > 1 and kwargs.get('do_metadata', False):
                    document_paragraphs.append(self.get_style_paragraph(block))
            elif type(block) is Table:
                table_tmp = []
                for j, row in enumerate(block.rows):
                    text += "|"
                    for k, cell in enumerate(row.cells):
                        text += cell.text + " | "
                        if len(cell.text) > 1 and kwargs.get('do_metadata', False):
                            document_paragraphs.append(self.get_style_paragraph(cell.paragraphs[0]))
                        if kwargs.get('do_tables', False):
                            table_tmp.append({
                                'text': cell.text,
                                'row': j,
                                'column': k,
                                'page': 0,
                            })
                    text += "\n"
                text += "\n"

                if kwargs.get('do_tables', False):
                    document_tables.append(table_tmp)

        extraction = {FULL_TEXT_KEY: text}

        return extraction, document_paragraphs, document_tables, document_lines

    def get_multiple_text(self, files: list, **kwargs: dict):
        """ Get text from multiple txt files

        :param files: (list) List of txt files to get text from
        :param kwargs: (dict) Additional parameters
        """
        raise NotImplementedError("This method is not implemented yet")

    def get_text_from_bytes(self, file: str):
        """ Get text from file

        :param file: File to get text from
        :return: (string) Text of the file
        """
        raise NotImplementedError("This method is not implemented yet")

    def extract_images(self, file: str, **kwargs: dict) -> list:
        """ Extract the images of the file and save them

        :param file: (str) File to extract the images from
        :param kwargs: (dict) Additional parameters
        :return: (list) List of images extracted
        """
        raise NotImplementedError("This method is not implemented yet")

    def extract_text(self, file: str, **kwargs: dict) -> Tuple[str, list, list, list]:
        """ Extract text from file and store it into a text file

        :param file: (str) File to extract text from
        :param kwargs: (dict) Additional parameters
        :return: (tuple) Filename, pages, paragraphs and tables extracted
        """
        extraction, paragraphs, tables, lines = self.get_text(file, **kwargs)
        text = extraction[0]['text']
        paragraphs = extraction[1]
        table = extraction[2]

        filename = os.path.splitext(file)[0] + ".txt"
        with open(filename, "+w") as f:
            for c in text:
                try:
                    f.write(c)
                except:
                    continue

        filename_pages = os.path.splitext(file)[0] + "_pag_0.txt"
        with open(filename_pages, "+w") as f:
            for c in text:
                try:
                    f.write(c)
                except:
                    continue

        filename_paragraphs = os.path.splitext(file)[0] + f"_paragraphs.json"
        with open(filename_paragraphs, "+w", encoding='utf8') as f:
            json.dump(paragraphs, f)
        filename_tables = os.path.splitext(file)[0] + f"_tables.json"
        with open(filename_tables, "+w", encoding='utf8') as f:
            json.dump(table, f)

        return filename, filename_pages, filename_paragraphs, filename_tables

    @staticmethod
    def write(text: str, filename: str):
        """ Write text into a file

        :param text: (str) Text to write
        :param filename: (str) Filename to write the text into
        """
        with open(filename, "w") as f:
            f.write(text)

    def extract_multiple_text(self, files: list, **kwargs: dict) -> Tuple[list, list, list, list]:
        """ Extract text from multiple files

        :param files: (list) List of files to extract text from
        :param kwargs: (dict) Additional parameters
        :return: (tuple) Filenames, pages, paragraphs and tables extracted
        """
        filenames = []
        filenames_pages = []
        filenames_paragraphs = []
        filenames_tables = []
        for file in files:
            filename_tmp, filenames_pages_tmp, filename_paragraphs_tmp, filename_tables_tmp = self.extract_text(file, **kwargs)
            filenames.append(filename_tmp)
            filenames_pages.append(filenames_pages_tmp)
            filenames_paragraphs.append(filename_paragraphs_tmp)
            filenames_tables.append(filename_tables_tmp)

        return filenames, filenames_pages, filenames_paragraphs, filenames_tables

    def get_number_pages(self, file: str) -> int:
        """ Get the number of pages of the file

        :param file: (str) File to get the number of pages from
        :return: (int) Number of pages of the file
        """
        try:
            pages = 1
            with open(file, 'rb') as f:
                zip_content = f.read()
            zip = zipfile.ZipFile(BytesIO(zip_content))
            with zip.open("docProps/app.xml") as xml:
                contenido = xml.read()
            root = et.fromstring(contenido)
            for element in root.iter():
                if "Pages" in element.tag:
                    pages = element.text
            n_pags = int(pages)
        except:
            n_pags = 1

        return n_pags


class PowerPointService(BaseFileService):
    TYPE_FILE = ["pptx"]

    def __init__(self):
        self.logger = logging.getLogger(__name__)

    @staticmethod
    def get_style_paragraph(shape: dict, page: int, width: float, height: float) -> dict:
        """ Get the style of a paragraph

        :param shape: (dict) Shape to get the style from
        :param page: (int) Page of the shape
        :param width: (float) Width of the slide
        :param height: (float) Height of the slide
        :return: (dict) Style of the paragraph
        """
        font = shape.text_frame.paragraphs[0].font
        style = {
            'text': shape.text,
            'bold': font.bold if font.bold else False,
            'italic': font.italic if font.italic else False,
            'underline': font.underline if font.underline else False,
            'font': font.name if font.name else "",
            'fontSize': font.size if font.size else "",
            'r0': shape.left.pt,
            'c0': shape.top.pt,
            'r1': shape.left.pt + shape.width.pt,
            'c1': shape.top.pt + shape.height.pt,
            'page': page,
            'slide_width': height,
            'slide_height': width
        }
        return style

    @staticmethod
    def get_style_cell(shape: dict, cell: dict, page: int, width: float, height: float) -> dict:
        """ Get the style of a cell

        :param shape: (dict) Shape to get the style from
        :param cell: (dict) Cell to get the style from
        :param page: (int) Page of the shape
        :param width: (float) Width of the slide
        :param height: (float) Height of the slide
        :return: (dict) Style of the cell
        """
        font = cell.text_frame.paragraphs[0].font
        style = {
            'text': cell.text,
            'bold': font.bold if font.bold else False,
            'italic': font.italic if font.italic else False,
            'underline': font.underline if font.underline else False,
            'font': font.name if font.name else "",
            'fontSize': font.size if font.size else "",
            'r0': shape.left.pt,
            'c0': shape.top.pt,
            'r1': shape.left.pt + shape.width.pt,
            'c1': shape.top.pt + shape.height.pt,
            'page': page,
            'slide_width': height,
            'slide_height': width
        }
        return style

    def get_text(self, file: str, **kwargs: dict) -> Tuple[dict, list, list, list]:
        """ Get the text from a file

        :param file: (str) File to get the text from
        :param kwargs: (dict) Additional parameters
        :return: (tuple) Text, paragraphs, tables and lines extracted
        """
        document_paragraphs = []
        document_tables = []
        document_lines = []
        extraction = {FULL_TEXT_KEY: ""}
        prs = Presentation(file)
        for i in range(self.get_number_pages(file)):
            text = ""
            for shape in prs.slides[i].shapes:
                if shape.has_text_frame:
                    text += shape.text
                    if len(shape.text) >= 1 and kwargs.get('do_metadata', False):
                        document_paragraphs.append(
                            self.get_style_paragraph(shape, i, prs.slide_width.pt, prs.slide_height.pt))
                if shape.has_table:
                    tbl = shape.table
                    table_tmp = []
                    for j in range(0, len(tbl.rows)):
                        text += "|"
                        for k in range(0, len(tbl.columns)):
                            cell = tbl.cell(j, k)
                            text += cell.text + " | "
                            if len(cell.text) >= 1 and kwargs.get('do_metadata', False):
                                document_paragraphs.append(
                                    self.get_style_cell(shape, cell, i, prs.slide_width.pt, prs.slide_height.pt))
                            if kwargs.get('do_tables', False):
                                table_tmp.append({
                                    'text': cell.text,
                                    'row': j,
                                    'column': k,
                                    'page': i,
                                    'slide_width': prs.slide_width.pt,
                                    'slide_height': prs.slide_height.pt,
                                    'r0': shape.left.pt,
                                    'c0': shape.top.pt,
                                    'r1': shape.left.pt + shape.width.pt,
                                    'c1': shape.top.pt + shape.height.pt
                                })
                        text += "\n"
                    text += "\n"
                    if kwargs.get('do_tables', False):
                        document_tables.append(table_tmp)
            extraction[f"pag_{i}"] = text
            text += "\n\n"
            extraction[FULL_TEXT_KEY] += text

        return extraction, document_paragraphs, document_tables, document_lines

    def get_multiple_text(self, files: list, **kwargs: dict):
        """ Get text from multiple txt files

        :param files: (list) List of txt files to get text from
        :param kwargs: (dict) Additional parameters
        """
        raise NotImplementedError("This method is not implemented yet")

    def get_text_from_bytes(self, file: str):
        """ Get text from file

        :param file: File to get text from
        :return: (string) Text of the file
        """
        raise NotImplementedError("This method is not implemented yet")

    def extract_images(self, file: str, **kwargs: dict) -> list:
        """ Extract the images of the file and save them

        :param file: (str) File to extract the images from
        :param kwargs: (dict) Additional parameters
        :return: (list) List of images extracted
        """
        raise NotImplementedError("This method is not implemented yet")

    def extract_text(self, file: str, **kwargs: dict) -> Tuple[str, list, list, list]:
        """ Extract text from file and store it into a text file

        :param file: (str) File to extract text from
        :param kwargs: (dict) Additional parameters
        :return: (tuple) Text, paragraphs, tables and lines extracted
        """
        extraction, paragraphs, tables, lines = self.get_text(file, **kwargs)
        filename = ""
        filenames_pages = []
        for key in extraction:
            fn = os.path.splitext(file)[0]
            if key == FULL_TEXT_KEY:
                fn += ".txt"
                filename = fn
            else:
                fn += f"_{key}.txt"
                filenames_pages.append(fn)
            with open(fn, "+w", encoding="utf8") as f:
                for c in extraction[key]:
                    try:
                        f.write(c)
                    except:
                        continue

        filename_paragraphs = os.path.splitext(file)[0] + f"_paragraphs.json"
        with open(filename_paragraphs, "+w", encoding="utf8") as f:
            json.dump(paragraphs, f)

        filename_tables = os.path.splitext(file)[0] + f"_tables.json"
        with open(filename_tables, "+w", encoding="utf8") as f:
            json.dump(tables, f)

        return filename, filenames_pages, filename_paragraphs, filename_tables

    def extract_multiple_text(self, files, **kwargs):
        """ Extract text from several files and store it into their corresponding text file

        :param files: Files to extract text from
        """
        filenames = []
        filenames_pages = []
        filenames_paragraphs = []
        filenames_tables = []
        for file in files:
            filename_tmp, filenames_pages_tmp, filename_paragraphs_tmp, filenames_tables_tmp = self.extract_text(file, **kwargs)
            filenames.append(filename_tmp)
            filenames_pages.append(filenames_pages_tmp)
            filenames_paragraphs.append(filename_paragraphs_tmp)
            filenames_tables.append(filenames_tables_tmp)
        return filenames, filenames_pages, filenames_paragraphs, filenames_tables

    def get_number_pages(self, file: str) -> int:
        """ Extract text from several files and store it into their corresponding text file

        :param file: (str) Files to extract text from
        :return: (int) Number of pages of the file
        """
        prs = Presentation(file)
        return len(prs.slides)


class ExcelService(BaseFileService):
    TYPE_FILE = ["xls", "xlsx"]

    def __init__(self):
        self.logger = logging.getLogger(__name__)

    @staticmethod
    def get_style_paragraph(cell: dict, page: int) -> dict:
        """" Get the style of the cell

        :param cell: (dict) Cell to get the style from
        :param page: (int) Page of the cell
        :return: (dict) Style of the cell
        """
        style = {
            'text': cell.value,
            'bold': cell.font.b if cell.font.b else False,
            'italic': cell.font.i if cell.font.i else False,
            'underline': cell.font.u if cell.font.u else False,
            'font': cell.font.name if cell.font.name else "",
            'fontSize': cell.font.size if cell.font.size else "",
            'coordinate': cell.coordinate if cell.coordinate else "",
            'r': cell.column - 1,
            'c': cell.row - 1,
            'page': page
        }
        return style

    def get_text_xlsx(self, file: str, **kwargs: dict) -> Tuple[dict, list, list, list]:
        """ Get text from xlsx file
        
        :param file: (str) File to get text from
        :param kwargs: (dict) Additional parameters
        :return: (tuple) Text, paragraphs, tables and lines extracted
        """
        extraction = {FULL_TEXT_KEY: ""}
        document_words = []
        document_tables = []
        document_lines = []

        wb = load_workbook(file)
        for i in range(self.get_number_pages(file)):
            text = ""
            table_tmp = []
            for j, row in enumerate(wb.worksheets[i].iter_rows()):
                text += "|"
                for k, cell in enumerate(row):
                    if str(cell.value) != "None":
                        text += " " + str(cell.value) + " |"
                        if kwargs.get('do_metadata', False):
                            document_words.append(self.get_style_paragraph(cell, i))
                        if kwargs.get('do_tables', False):
                            table_tmp.append({
                                'text': cell.value,
                                'row': j,
                                'column': k,
                                'page': i,
                            })
                text += "\n"

            extraction[f"pag_{i}"] = text
            text += "\n\n"
            extraction[FULL_TEXT_KEY] += text

            if kwargs.get('do_tables', False):
                document_tables.append(table_tmp)

        return extraction, document_words, document_tables, document_lines

    def get_text_xls(self, file: str, **kwargs: dict) -> Tuple[dict, list, list, list]:
        """ Get text from xls file

        :param file: (str) File to get text from
        :param kwargs: (dict) Additional parameters
        :return: (tuple) Text, paragraphs, tables and lines extracted
        """
        extraction = {FULL_TEXT_KEY: ""}
        document_words = []
        document_tables = []
        document_lines = []

        workbook = xlrd.open_workbook(file)
        for i in range(self.get_number_pages(file)):
            row_count = workbook.sheets()[i].nrows
            col_count = workbook.sheets()[i].ncols
            text = ""
            table_tmp = []
            for cur_row in range(0, row_count):
                text += "|"
                for cur_col in range(0, col_count):
                    cell = workbook.sheets()[i].cell(cur_row, cur_col)
                    if str(cell.value) != "None":
                        text += " " + str(cell.value) + " |"
                        if kwargs.get('do_tables', False):
                            table_tmp.append({
                                'text': cell.value,
                                'row': cur_row,
                                'column': cur_col,
                                'page': i,
                            })
                text += "\n"

            extraction[f"pag_{i}"] = text
            text += "\n\n"
            extraction[FULL_TEXT_KEY] += text

            if kwargs.get('do_tables', False):
                document_tables.append(table_tmp)

        return extraction, document_words, document_tables, document_lines

    def get_text(self, file: str, **kwargs: dict):
        """ Get text from txt fle

        :param file: txt file to get text from
        :return: (String) Text of the txt file
        """
        extraction, pages_cells, pages_tables, lines = {}, [], [], []
        type = self.get_type(file)
        if type == "xls":
            extraction, pages_cells, pages_tables, lines = self.get_text_xls(file, **kwargs)
        elif type == "xlsx":
            name, ext = os.path.splitext(file)
            if ext != ".xlsx":
                new_file = name + ".xlsx"
                copy2(file, new_file)
            else:
                new_file = file
            extraction, pages_cells, pages_tables, lines = self.get_text_xlsx(new_file, **kwargs)

        return extraction, pages_cells, pages_tables, lines

    def get_multiple_text(self, files: list, **kwargs: dict):
        """ Get text from multiple txt files

        :param files: (list) List of txt files to get text from
        :param kwargs: (dict) Additional parameters
        """
        raise NotImplementedError("This method is not implemented yet")

    def get_text_from_bytes(self, file: str):
        """ Get text from file

        :param file: File to get text from
        :return: (string) Text of the file
        """
        raise NotImplementedError("This method is not implemented yet")

    def extract_images(self, file: str, **kwargs: dict) -> list:
        """ Extract the images of the file and save them

        :param file: (str) File to extract the images from
        :param kwargs: (dict) Additional parameters
        :return: (list) List of images extracted
        """
        raise NotImplementedError("This method is not implemented yet")

    def extract_text(self, file: str, **kwargs: dict) -> Tuple[str, list, list, list]:
        """ Extract text from file and store it into a text file

        :param file: (str) File to extract text from
        :param kwargs: (dict) Additional parameters
        :return: (tuple) Text, paragraphs, tables and lines extracted
        """
        extraction, words, tables, lines = self.get_text(file, **kwargs)
        filename = ""
        filenames_pages = []

        for key in extraction:
            fn = os.path.splitext(file)[0]
            if key == FULL_TEXT_KEY:
                fn += ".txt"
                filename = fn
            else:
                fn += f"_{key}.txt"
                filenames_pages.append(fn)

            with open(fn, "+w", encoding="utf8") as f:
                for c in extraction[key]:
                    try:
                        f.write(c)
                    except:
                        continue

        filename_words = os.path.splitext(file)[0] + f"_words.json"
        with open(filename_words, "+w", encoding="utf8") as f:
            json.dump(words, f)

            filename_tables = os.path.splitext(file)[0] + f"_tables.json"
            with open(filename_tables, "+w", encoding="utf8") as fo:
                json.dump(tables, fo)

        return filename, filenames_pages, filename_words, filename_tables

    def extract_multiple_text(self, files: list, **kwargs: dict) -> Tuple[list, list, list, list]:
        """ Extract text from several files and store it into their corresponding text file

        :param files: (list) Files to extract text from
        :param kwargs: (dict) Additional parameters
        :return: (tuple) List of filenames of the text files
        """
        filenames = []
        filenames_pages = []
        filenames_words = []
        filenames_tables = []

        for file in files:
            filename_tmp, filenames_pages_tmp, filename_words_tmp, filename_tables_tmp = self.extract_text(file, **kwargs)
            filenames.append(filename_tmp)
            filenames_pages.append(filenames_pages_tmp)
            filenames_words.append(filename_words_tmp)
            filenames_tables.append(filename_tables_tmp)

        return filenames, filenames_pages, filenames_words, filenames_tables

    def get_number_pages(self, file: str) -> int:
        """ Extract text from several files and store it into their corresponding text file

        :param file: Files to extract text from
        :return: (int) Number of pages of the file
        """
        type = self.get_type(file)
        if type == "xls":
            workbook = xlrd.open_workbook(file)
            return len(workbook.sheets())
        elif type == "xlsx":
            name, ext = os.path.splitext(file)
            if ext != ".xlsx":
                new_file = name + ".xlsx"
                copy2(file, new_file)
            else:
                new_file = file
            wb = load_workbook(new_file)
            return len(wb.worksheets)

    @staticmethod
    def get_type(filename: str) -> str:
        """ Get the type of the file

        :param filename: Name of the file to get type of
        :return: (str) Type of the file
        """
        file_formats_map = {
            'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet': "xlsx",
            'application/vnd.ms-excel': "xls"
        }
        mime = magic.from_buffer(open(filename, "rb").read(1024*1024), mime=True)
        return file_formats_map.get(mime, filename.split(".")[-1])
